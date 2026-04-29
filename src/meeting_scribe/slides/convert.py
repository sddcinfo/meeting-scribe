"""Pure conversion functions for PPTX processing.

These functions run INSIDE the sandboxed worker container. They have no
dependency on the meeting-scribe server, FastAPI, or any backend.
"""

from __future__ import annotations

import json
import logging
import shutil
import subprocess
import uuid
import zipfile
from collections.abc import Callable
from pathlib import Path

from meeting_scribe.slides.models import SlideText, TextRun, ValidationResult

logger = logging.getLogger(__name__)

# Trusted local environment — minimal limits
MAX_SLIDES = 200
MAX_OUTPUT_MB = 500
LIBREOFFICE_TIMEOUT = 120  # seconds per conversion

ProgressCb = Callable[[int, int], None]  # (slide_index_0based, total_slides)

# ── Validation ───────────────────────────────────────────────


def validate_pptx_contents(pptx_path: Path) -> ValidationResult:
    """Validate a PPTX file — just check it's parseable and count slides."""
    if not zipfile.is_zipfile(pptx_path):
        return ValidationResult(valid=False, error="Not a valid ZIP/PPTX file")

    try:
        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        slide_count = len(prs.slides)

        if slide_count == 0:
            return ValidationResult(valid=False, error="Presentation has no slides")

        return ValidationResult(
            valid=True,
            slide_count=slide_count,
            slide_width=int(prs.slide_width) if prs.slide_width else 0,
            slide_height=int(prs.slide_height) if prs.slide_height else 0,
        )

    except Exception as exc:
        return ValidationResult(valid=False, error=f"Failed to parse PPTX: {exc}")


# ── Text extraction ──────────────────────────────────────────


def _iter_shapes_recursive(shapes):
    """Yield every shape under ``shapes``, descending into GROUP shapes.

    python-pptx's default ``slide.shapes`` only iterates the top-level
    container — anything inside a group is invisible. This walker fixes the
    most common "text was on the slide but didn't appear in extract" bug.
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception:  # pragma: no cover — pptx not installed
        MSO_SHAPE_TYPE = None  # type: ignore[assignment]

    for shape in shapes:
        is_group = (
            MSO_SHAPE_TYPE is not None
            and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP
        )
        if is_group:
            yield from _iter_shapes_recursive(shape.shapes)
        else:
            yield shape


def extract_text_from_pptx(pptx_path: Path) -> list[SlideText]:
    """Extract all text runs from a PPTX with formatting metadata.

    Returns one SlideText per slide, each containing TextRun objects
    with stable IDs for round-trip translation.

    Walks group shapes recursively. ``shape_id`` is unique per slide in
    OOXML regardless of nesting depth, so the existing run-id format keeps
    its uniqueness contract.
    """
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slides: list[SlideText] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_text = SlideText(index=slide_idx)

        for shape in _iter_shapes_recursive(slide.shapes):
            if shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    for run_idx, run in enumerate(para.runs):
                        text = run.text.strip()
                        if not text:
                            continue

                        run_id = f"s{slide_idx}_sh{shape.shape_id}_p{para_idx}_r{run_idx}"
                        slide_text.runs.append(
                            TextRun(
                                id=run_id,
                                slide_index=slide_idx,
                                shape_id=shape.shape_id,
                                para_index=para_idx,
                                run_index=run_idx,
                                text=text,
                                font_name=run.font.name,
                                font_size=run.font.size,
                                bold=run.font.bold,
                                italic=run.font.italic,
                            )
                        )

            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        for para_idx, para in enumerate(cell.text_frame.paragraphs):
                            for run_idx, run in enumerate(para.runs):
                                text = run.text.strip()
                                if not text:
                                    continue

                                run_id = (
                                    f"s{slide_idx}_sh{shape.shape_id}"
                                    f"_t{row_idx}x{col_idx}_p{para_idx}_r{run_idx}"
                                )
                                slide_text.runs.append(
                                    TextRun(
                                        id=run_id,
                                        slide_index=slide_idx,
                                        shape_id=shape.shape_id,
                                        para_index=para_idx,
                                        run_index=run_idx,
                                        text=text,
                                        font_name=run.font.name,
                                        font_size=run.font.size,
                                        bold=run.font.bold,
                                        italic=run.font.italic,
                                    )
                                )

        slides.append(slide_text)

    # Secondary pass: extract SmartArt text (python-pptx misses these)
    try:
        from meeting_scribe.slides.smartart import extract_smartart_text

        smartart = extract_smartart_text(pptx_path)
        for slide_idx, texts in smartart.items():
            if slide_idx < len(slides):
                for i, text in enumerate(texts):
                    if text.strip():
                        run_id = f"s{slide_idx}_smartart_{i}"
                        slides[slide_idx].runs.append(
                            TextRun(
                                id=run_id,
                                slide_index=slide_idx,
                                shape_id=0,
                                para_index=0,
                                run_index=i,
                                text=text.strip(),
                            )
                        )
    except Exception:
        logger.debug("SmartArt extraction skipped", exc_info=True)

    return slides


# ── Rendering ────────────────────────────────────────────────


def extract_embedded_thumbnail(pptx_path: Path, dest_png: Path) -> bool:
    """Pull the embedded thumbnail (docProps/thumbnail.jpeg) out of a PPTX.

    PowerPoint and Keynote always write a small (~256x144) preview JPEG to
    ``docProps/thumbnail.jpeg`` when saving. Reading it is a sub-50ms
    zip-read with no rendering needed, so it makes a great instant
    placeholder for the first slide while the real LibreOffice render
    runs in the background.

    Converts JPEG → PNG via Pillow so the on-disk file format matches its
    extension (and the HTTP endpoint can advertise image/png honestly).

    Returns True if a thumbnail was extracted and saved as ``dest_png``.
    Falsy if the file doesn't have one (often the case for python-pptx
    or Google Slides exports).
    """
    candidates = (
        "docProps/thumbnail.jpeg",
        "docProps/thumbnail.jpg",
        "docProps/thumbnail.png",
    )
    try:
        with zipfile.ZipFile(pptx_path) as zf:
            names = set(zf.namelist())
            for cand in candidates:
                if cand not in names:
                    continue
                data = zf.read(cand)
                if not data:
                    continue
                dest_png.parent.mkdir(parents=True, exist_ok=True)
                tmp = dest_png.with_suffix(dest_png.suffix + ".tmp")
                if cand.endswith(".png"):
                    tmp.write_bytes(data)
                else:
                    # Re-encode JPEG → PNG so MIME type matches the extension
                    try:
                        from io import BytesIO

                        from PIL import Image

                        img = Image.open(BytesIO(data)).convert("RGB")
                        img.save(tmp, format="PNG", optimize=False)
                    except Exception:
                        # Pillow unavailable or decode failed — fall back to
                        # raw write; browsers sniff magic bytes regardless of
                        # the .png extension and render JPEG content fine.
                        tmp.write_bytes(data)
                tmp.replace(dest_png)
                return True
    except zipfile.BadZipFile, OSError:
        logger.debug("No usable embedded thumbnail in %s", pptx_path)
    return False


_UNOSERVER_HOST = "127.0.0.1"
_UNOSERVER_PORT_ENV = "MEETING_SCRIBE_UNOSERVER_PORT"
_UNOSERVER_DEFAULT_PORT = 2003


def _unoserver_endpoint() -> tuple[str, int] | None:
    """Return (host, port) of a reachable unoserver, or None."""
    import os
    import socket

    port_env = os.environ.get(_UNOSERVER_PORT_ENV, str(_UNOSERVER_DEFAULT_PORT))
    try:
        port = int(port_env)
    except ValueError:
        return None
    try:
        with socket.create_connection((_UNOSERVER_HOST, port), timeout=0.2):
            return (_UNOSERVER_HOST, port)
    except TimeoutError, OSError:
        return None


def _libreoffice_pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    """Convert PPTX → PDF, preferring a warm unoserver if reachable.

    The unoserver path keeps a single LibreOffice process resident, so
    subsequent conversions skip the ~5s cold-start. For express renders
    (single-slide minimal PPTX) this drops the wall-clock from ~3-5s to
    well under 1s. Cold start of unoserver itself happens once at boot.

    Falls back to the subprocess `libreoffice --convert-to` path if
    unoserver is not running. No new install is required for the fallback;
    the daemon path is opt-in.
    """
    out_dir.mkdir(parents=True, exist_ok=True)
    dest_pdf = out_dir / (pptx_path.stem + ".pdf")

    endpoint = _unoserver_endpoint()
    if endpoint is not None:
        try:
            host, port = endpoint
            subprocess.run(
                [
                    "unoconvert",
                    "--host",
                    host,
                    "--port",
                    str(port),
                    "--convert-to",
                    "pdf",
                    str(pptx_path),
                    str(dest_pdf),
                ],
                capture_output=True,
                text=True,
                timeout=LIBREOFFICE_TIMEOUT,
                check=True,
            )
            if dest_pdf.exists():
                return dest_pdf
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError) as exc:
            # FileNotFoundError = unoconvert binary missing; fall through
            logger.debug("unoserver path failed (%s); falling back to subprocess", exc)

    # Subprocess fallback. Each invocation gets its OWN UserInstallation
    # profile dir — without this, two concurrent libreoffice processes
    # (express render + bulk render of the same upload) fight over
    # ~/.config/libreoffice/4/user/.~lock and one is silently killed
    # (returncode!=0, stderr empty). `--nolockcheck` doesn't bypass this
    # reliably; isolating the profile does.
    import tempfile
    import uuid as _uuid

    user_dir = Path(tempfile.gettempdir()) / f"lo_user_{_uuid.uuid4().hex[:12]}"
    user_dir.mkdir(parents=True, exist_ok=True)
    user_install_uri = "file://" + str(user_dir)
    try:
        # No --infilter: let LibreOffice auto-detect the OOXML variant from
        # the content type. The previous hardcoded
        # `impress_MS_PowerPoint_2007_XML` worked for .pptx but rejected
        # template (.potx) and slideshow (.ppsx) content types — even
        # when renamed to .pptx, the embedded content-type header still
        # identifies them as templates and the wrong filter silently
        # produces no PDF (rc!=0, empty stderr — same failure mode that
        # crashed the user's last upload).
        result = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--norestore",
                "--nolockcheck",
                "--nologo",
                "-env:UserInstallation=" + user_install_uri,
                "--convert-to",
                "pdf",
                "--outdir",
                str(out_dir),
                str(pptx_path),
            ],
            capture_output=True,
            text=True,
            timeout=LIBREOFFICE_TIMEOUT,
        )
        if result.returncode != 0:
            # Empty stderr usually = killed by signal (OOM, profile lock
            # collision). Surface stdout too so callers can see what little
            # libreoffice did emit before dying.
            tail = (result.stderr or "").strip() or (result.stdout or "").strip()
            logger.error(
                "LibreOffice failed (rc=%d): %s",
                result.returncode,
                tail or "(no output — process likely killed by signal)",
            )
            raise RuntimeError(
                f"LibreOffice conversion failed (rc={result.returncode}): "
                f"{tail[:500] or 'process killed (no output)'}"
            )

        pdf_files = sorted(out_dir.glob("*.pdf"))
        if not pdf_files:
            raise RuntimeError("LibreOffice produced no PDF output")
        return pdf_files[0]
    finally:
        shutil.rmtree(user_dir, ignore_errors=True)


def _pdf_page_count(pdf_path: Path) -> int:
    """Return the number of pages in a PDF using pdfinfo."""
    result = subprocess.run(
        ["pdfinfo", str(pdf_path)],
        capture_output=True,
        text=True,
        check=True,
        timeout=15,
    )
    for line in result.stdout.splitlines():
        if line.startswith("Pages:"):
            return int(line.split(":", 1)[1].strip())
    raise RuntimeError("pdfinfo did not report page count")


def _pdftoppm_single_page(pdf_path: Path, page_num: int, dest_png: Path) -> None:
    """Render a single PDF page to a PNG at the given destination path.

    The temp prefix carries a random token so concurrent calls writing
    to the same ``dest_png`` (for example, the express-first-slide
    renderer and the bulk-originals renderer both targeting
    ``original/slide_001.png``) cannot race on the same `.tmp.png`
    file.  Without the token, a rename by one caller could orphan the
    other caller's in-progress pdftoppm output and surface as
    ``RuntimeError: pdftoppm produced no output``, which aborts the
    bulk render partway through and leaves ``original/`` with only the
    first slide.
    """
    # pdftoppm with -singlefile appends .png to the prefix we pass, so
    # we strip the dest suffix and let pdftoppm add it back.  The
    # uuid keeps two concurrent renders from clobbering each other's
    # temp file.
    tmp_prefix = dest_png.with_suffix(f".{uuid.uuid4().hex}.tmp")
    subprocess.run(
        [
            "pdftoppm",
            "-png",
            "-r",
            "300",
            "-f",
            str(page_num),
            "-l",
            str(page_num),
            "-singlefile",
            str(pdf_path),
            str(tmp_prefix),
        ],
        capture_output=True,
        text=True,
        timeout=LIBREOFFICE_TIMEOUT,
        check=True,
    )
    # Two output-path conventions across pdftoppm versions:
    #   * mainline: strips the final extension from the prefix and
    #     appends ".png" → ``tmp_prefix.with_suffix(".png")``;
    #   * older:    appends ".png" literally to whatever we pass →
    #     ``Path(str(tmp_prefix) + ".png")``.
    # Try both so we are portable across packaging ages.
    candidates = [
        tmp_prefix.with_suffix(".png"),
        Path(str(tmp_prefix) + ".png"),
    ]
    tmp_png = next((c for c in candidates if c.exists()), None)
    if tmp_png is None:
        raise RuntimeError(
            f"pdftoppm produced no output for page {page_num} "
            f"(tried {[str(c) for c in candidates]})"
        )
    tmp_png.replace(dest_png)


def convert_pptx_to_images(
    pptx_path: Path,
    output_dir: Path,
    progress_cb: ProgressCb | None = None,
) -> int:
    """Convert PPTX to PDF + per-slide PNGs via LibreOffice headless.

    Renders slide 1 FIRST (priority pass) so popout viewers can display the
    first slide as soon as the LibreOffice PDF conversion completes. Remaining
    slides render sequentially after. ``progress_cb(slide_idx_0based, total)``
    fires after each PNG is finalized on disk (atomic rename).

    The PDF is kept as ``original.pdf`` for the source-side viewer (lossless
    vector rendering via Chrome's built-in PDF viewer). Per-slide PNGs are
    generated at 300 DPI as fallback and for translated slides.

    Returns the number of slides rendered.
    """
    output_dir.mkdir(parents=True, exist_ok=True)

    tmp_dir = Path("/tmp/lo_render")
    tmp_dir.mkdir(parents=True, exist_ok=True)

    try:
        pdf_path = _libreoffice_pptx_to_pdf(pptx_path, tmp_dir)

        # Keep the PDF for source-side viewing (lossless quality)
        shutil.copy2(str(pdf_path), str(output_dir / "original.pdf"))

        total = _pdf_page_count(pdf_path)
        if total == 0:
            raise RuntimeError("PDF reports zero pages")

        running_size = 0
        max_output_bytes = MAX_OUTPUT_MB * 1024 * 1024

        def _render_one(page_num: int) -> None:
            nonlocal running_size
            dest = output_dir / f"slide_{page_num:03d}.png"
            _pdftoppm_single_page(pdf_path, page_num, dest)
            running_size += dest.stat().st_size
            if running_size > max_output_bytes:
                raise RuntimeError(
                    f"Output size {running_size} bytes exceeds limit ({MAX_OUTPUT_MB} MB)"
                )
            if progress_cb is not None:
                try:
                    progress_cb(page_num - 1, total)
                except Exception:
                    logger.debug("progress_cb raised", exc_info=True)

        # Priority pass: slide 1 first so the popout has something to show
        _render_one(1)
        # Remaining slides in order
        for page_num in range(2, total + 1):
            _render_one(page_num)

        return total

    finally:
        # Clean up temp rendering directory
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ── Text reinsertion ─────────────────────────────────────────


def reinsert_translated_text(
    pptx_path: Path,
    translated_runs: list[dict],  # [{"id": "...", "translated": "..."}]
    output_path: Path,
) -> None:
    """Replace text in PPTX shapes with translated text, preserving formatting.

    Each entry in translated_runs maps an id (from extract_text_from_pptx)
    to its translated text. Formatting (font, size, color, bold, italic)
    is preserved at the run level.
    """
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE

    prs = Presentation(str(pptx_path))
    tr_map = {r["id"]: r["translated"] for r in translated_runs}

    for slide_idx, slide in enumerate(prs.slides):
        for shape in _iter_shapes_recursive(slide.shapes):
            _replace_shape_text(shape, slide_idx, tr_map)

            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        _replace_cell_text(
                            cell,
                            slide_idx,
                            shape.shape_id,
                            row_idx,
                            col_idx,
                            tr_map,
                        )

        # Set auto-fit on all text frames to handle overflow (group-aware)
        for shape in _iter_shapes_recursive(slide.shapes):
            if shape.has_text_frame:
                try:
                    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception:
                    pass  # Some shapes don't support auto_size

    prs.save(str(output_path))


def _replace_shape_text(
    shape,
    slide_idx: int,
    tr_map: dict[str, str],
) -> None:
    """Replace text runs in a shape using the translation map."""
    if not shape.has_text_frame:
        return

    for para_idx, para in enumerate(shape.text_frame.paragraphs):
        for run_idx, run in enumerate(para.runs):
            run_id = f"s{slide_idx}_sh{shape.shape_id}_p{para_idx}_r{run_idx}"
            if run_id in tr_map:
                original_len = len(run.text)
                translated = tr_map[run_id]
                run.text = translated

                # Post-check: if translated text is much longer, reduce font
                if run.font.size and len(translated) > original_len * 1.5:
                    scale = original_len / max(len(translated), 1)
                    run.font.size = int(run.font.size * max(scale, 0.5))


def _replace_cell_text(
    cell,
    slide_idx: int,
    shape_id: int,
    row_idx: int,
    col_idx: int,
    tr_map: dict[str, str],
) -> None:
    """Replace text runs in a table cell."""
    for para_idx, para in enumerate(cell.text_frame.paragraphs):
        for run_idx, run in enumerate(para.runs):
            run_id = f"s{slide_idx}_sh{shape_id}_t{row_idx}x{col_idx}_p{para_idx}_r{run_idx}"
            if run_id in tr_map:
                original_len = len(run.text)
                translated = tr_map[run_id]
                run.text = translated

                if run.font.size and len(translated) > original_len * 1.5:
                    scale = original_len / max(len(translated), 1)
                    run.font.size = int(run.font.size * max(scale, 0.5))


def render_translated_to_images(
    pptx_path: Path,
    output_dir: Path,
    progress_cb: ProgressCb | None = None,
) -> int:
    """Render translated PPTX to PNG images. Same as convert_pptx_to_images."""
    return convert_pptx_to_images(pptx_path, output_dir, progress_cb=progress_cb)


def build_minimal_pptx(source_pptx: Path, keep_indices_0based: list[int], dest_pptx: Path) -> int:
    """Build a stripped-down PPTX containing only the requested slide indices.

    Removes other slide entries from ``ppt/presentation.xml``'s ``<sldIdLst>``
    and the matching relationships in ``ppt/_rels/presentation.xml.rels``,
    then re-zips excluding the orphaned ``ppt/slides/slideN.xml`` parts and
    their _rels (otherwise LibreOffice still parses them on open and we lose
    most of the speed benefit).

    Theme/master/layout parts are owned by the presentation, not slide parts,
    so they survive cleanly. Returns the number of slides kept.

    This is the express-lane workhorse — it lets us run LibreOffice on a
    1-slide deck (~3-5s cold) instead of the full deck (~25-30s for 50
    slides), letting the first translated slide reach the popout much sooner.
    """
    from xml.etree import ElementTree as ET

    keep_set = set(keep_indices_0based)
    if not keep_set:
        raise ValueError("keep_indices_0based must contain at least one index")

    NS = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
    }
    for prefix, uri in NS.items():
        ET.register_namespace("" if prefix == "rel" else prefix, uri)

    dest_pptx.parent.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(source_pptx, "r") as zin:
        names = list(zin.namelist())

        pres_xml = zin.read("ppt/presentation.xml")
        pres_rels_xml = zin.read("ppt/_rels/presentation.xml.rels")

        pres_root = ET.fromstring(pres_xml)
        sldIdLst = pres_root.find(f"{{{NS['p']}}}sldIdLst")
        if sldIdLst is None:
            raise RuntimeError("presentation.xml has no sldIdLst")

        all_sldids = list(sldIdLst)
        # Map slide-index → relationship id (rId) → target part name
        rels_root = ET.fromstring(pres_rels_xml)
        rel_target_by_id: dict[str, str] = {}
        for rel in rels_root.findall(f"{{{NS['rel']}}}Relationship"):
            rid = rel.get("Id")
            target = rel.get("Target", "")
            if rid:
                rel_target_by_id[rid] = target

        keep_rids: set[str] = set()
        drop_targets: set[str] = set()
        kept_count = 0
        for idx, sldid in enumerate(all_sldids):
            rid = sldid.get(f"{{{NS['r']}}}id")
            if idx in keep_set:
                if rid:
                    keep_rids.add(rid)
                kept_count += 1
            else:
                if rid:
                    target = rel_target_by_id.get(rid, "")
                    # Targets look like "slides/slide5.xml" — normalize to
                    # the full archive path so we can drop it cleanly.
                    if target:
                        if target.startswith("/"):
                            full = target.lstrip("/")
                        else:
                            full = "ppt/" + target
                        drop_targets.add(full)
                        drop_targets.add(
                            full.replace(".xml", ".xml.rels").replace(
                                "ppt/slides/", "ppt/slides/_rels/"
                            )
                        )
                # Remove the sldId entry from the presentation
                sldIdLst.remove(sldid)

        # Strip dropped slide rels from presentation.xml.rels
        for rel in list(rels_root.findall(f"{{{NS['rel']}}}Relationship")):
            rid = rel.get("Id")
            target = rel.get("Target", "")
            target_kind = ""
            if "slides/slide" in target and target.endswith(".xml"):
                target_kind = "slide"
            if target_kind == "slide" and rid not in keep_rids:
                rels_root.remove(rel)

        new_pres_xml = ET.tostring(pres_root, xml_declaration=True, encoding="UTF-8")
        new_rels_xml = ET.tostring(rels_root, xml_declaration=True, encoding="UTF-8")

        with zipfile.ZipFile(dest_pptx, "w", zipfile.ZIP_DEFLATED) as zout:
            for name in names:
                # Skip orphaned slide parts entirely so LO doesn't parse them
                if name in drop_targets:
                    continue
                # Also skip notes slides + comments tied to dropped slides
                # (best-effort by filename pattern; failures are non-fatal).
                base = name.rsplit("/", 1)[-1]
                if name.startswith("ppt/notesSlides/") and base.startswith("notesSlide"):
                    # Notes parts are usually 1:1 with slides — keeping them
                    # all is harmless; LO ignores notes during PNG render.
                    pass
                if name == "ppt/presentation.xml":
                    zout.writestr(name, new_pres_xml)
                elif name == "ppt/_rels/presentation.xml.rels":
                    zout.writestr(name, new_rels_xml)
                else:
                    zout.writestr(name, zin.read(name))

    return kept_count


def render_first_slide_fast(source_pptx: Path, dest_png: Path, work_dir: Path) -> bool:
    """Render JUST slide 1 of a PPTX as fast as possible.

    Strips the deck down to a 1-slide minimal PPTX, then runs LibreOffice
    on that. Drops cold render from ~25-30s (50-slide deck) to ~3-5s (1
    slide). Far faster than the priority pdftoppm trick which still pays
    LibreOffice's full-deck parse cost.

    Returns True on success. Failure is non-fatal — caller falls back to
    the existing full-deck render path.
    """
    work_dir.mkdir(parents=True, exist_ok=True)
    minimal = work_dir / "first_slide_only.pptx"
    try:
        build_minimal_pptx(source_pptx, [0], minimal)
        pdf_path = _libreoffice_pptx_to_pdf(minimal, work_dir)
        _pdftoppm_single_page(pdf_path, 1, dest_png)
        return True
    except Exception:
        logger.exception("render_first_slide_fast failed")
        return False
    finally:
        for f in work_dir.glob("*.pdf"):
            f.unlink(missing_ok=True)
        minimal.unlink(missing_ok=True)


def render_partial_translated(
    source_pptx: Path,
    translations: list[dict],
    slide_indices_0based: list[int],
    output_dir: Path,
    work_dir: Path,
) -> list[int]:
    """Express-lane render: produce translated PNGs for ONLY the given slides.

    Used to make the first 1-2 translated slides visible quickly, before the
    bulk reinsert + render of the entire deck completes.

    Pipeline:
      1. Reinsert translations into a copy of the source PPTX (full deck).
      2. Strip every slide except the targeted indices, producing a tiny
         minimal PPTX (drops LibreOffice parse cost from ~25s for a 50-slide
         deck to ~3-5s for a 1-slide one).
      3. Run LibreOffice on the minimal deck → PDF → per-page PNG.
      4. Save each PNG into ``output_dir/slide_NNN.png`` using the ORIGINAL
         deck's slide numbering (not the minimal-deck position).

    The bulk render later overwrites these PNGs with its final output.
    Returns the 0-based ORIGINAL indices that were successfully rendered.
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    work_dir.mkdir(parents=True, exist_ok=True)

    if not slide_indices_0based:
        return []

    partial_pptx = work_dir / "partial_translated.pptx"
    minimal_pptx = work_dir / "express_minimal.pptx"
    reinsert_translated_text(source_pptx, translations, partial_pptx)

    try:
        # Strip down to just the slides we need so LO has very little to parse.
        # Sort + dedupe so the kept slides land in deterministic deck order.
        targeted = sorted(set(slide_indices_0based))
        try:
            build_minimal_pptx(partial_pptx, targeted, minimal_pptx)
            render_input = minimal_pptx
        except Exception:
            # If splitting fails for any reason (unusual PPTX, locked rels,
            # etc.), fall back to rendering the full reinserted deck. We
            # still get correct output, just not the speed win.
            logger.exception("build_minimal_pptx failed; rendering full deck")
            render_input = partial_pptx

        pdf_path = _libreoffice_pptx_to_pdf(render_input, work_dir)
        total_pages = _pdf_page_count(pdf_path)

        rendered: list[int] = []
        if render_input is minimal_pptx:
            # Minimal PPTX has slides in `targeted` order — page N of the PDF
            # corresponds to targeted[N-1] in the original deck.
            for new_page_num, orig_idx0 in enumerate(targeted, start=1):
                if new_page_num > total_pages:
                    break
                dest = output_dir / f"slide_{orig_idx0 + 1:03d}.png"
                _pdftoppm_single_page(pdf_path, new_page_num, dest)
                rendered.append(orig_idx0)
        else:
            # Full-deck fallback: page numbers match the original indices.
            for idx0 in targeted:
                page_num = idx0 + 1
                if page_num < 1 or page_num > total_pages:
                    continue
                dest = output_dir / f"slide_{page_num:03d}.png"
                _pdftoppm_single_page(pdf_path, page_num, dest)
                rendered.append(idx0)
        return rendered
    finally:
        # Best-effort cleanup of intermediate artefacts
        for f in work_dir.glob("*.pdf"):
            f.unlink(missing_ok=True)
        partial_pptx.unlink(missing_ok=True)
        minimal_pptx.unlink(missing_ok=True)


# ── Serialization helpers ────────────────────────────────────


def slides_to_json(slides: list[SlideText]) -> list[dict]:
    """Serialize extracted text to JSON-friendly format."""
    return [
        {
            "index": s.index,
            "runs": [
                {
                    "id": r.id,
                    "slide_index": r.slide_index,
                    "shape_id": r.shape_id,
                    "text": r.text,
                    "font_name": r.font_name,
                    "font_size": r.font_size,
                    "bold": r.bold,
                    "italic": r.italic,
                }
                for r in s.runs
            ],
        }
        for s in slides
    ]


def write_text_extract(slides: list[SlideText], output_path: Path) -> None:
    """Write extracted text to JSON file."""
    data = slides_to_json(slides)
    output_path.write_text(json.dumps(data, ensure_ascii=False, indent=2))


# ── Language detection ───────────────────────────────────────


def _script_counts(slides: list[SlideText]) -> dict[str, int]:
    """Count characters by Unicode script across all slide runs.

    Ignores digits, punctuation, whitespace, and symbols — they would
    otherwise dilute CJK ratios on real-world decks heavy with dates,
    percentages, table numbers, and brand names.
    """
    counts = {"latin": 0, "hiragana": 0, "katakana": 0, "han": 0, "hangul": 0}
    for slide in slides:
        for run in slide.runs:
            for ch in run.text:
                cp = ord(ch)
                if 0x3040 <= cp <= 0x309F:
                    counts["hiragana"] += 1
                elif 0x30A0 <= cp <= 0x30FF or 0x31F0 <= cp <= 0x31FF:
                    counts["katakana"] += 1
                elif (
                    0x4E00 <= cp <= 0x9FFF  # CJK Unified Ideographs
                    or 0x3400 <= cp <= 0x4DBF  # CJK Extension A
                    or 0xF900 <= cp <= 0xFAFF  # CJK Compatibility
                    or 0x20000 <= cp <= 0x2A6DF  # CJK Extension B
                ):
                    counts["han"] += 1
                elif 0xAC00 <= cp <= 0xD7AF or 0x1100 <= cp <= 0x11FF:
                    counts["hangul"] += 1
                elif ch.isascii() and ch.isalpha():
                    counts["latin"] += 1
    return counts


def _detect_via_lingua(slides: list[SlideText]) -> str | None:
    """Use lingua-language-detector if installed; else None.

    lingua is the most accurate option for short / mixed text and reports
    per-language confidences directly, but adds a ~170 MB wheel. Treat as
    an opt-in upgrade — when not installed we fall back to the script-
    aware heuristic in detect_slide_language().
    """
    try:
        from lingua import Language, LanguageDetectorBuilder  # type: ignore[import-not-found]
    except Exception:
        return None

    corpus = " ".join(
        run.text.strip() for slide in slides for run in slide.runs if run.text and run.text.strip()
    ).strip()
    if not corpus:
        return None

    try:
        detector = LanguageDetectorBuilder.from_languages(
            Language.ENGLISH, Language.JAPANESE, Language.CHINESE, Language.KOREAN
        ).build()
        confs = detector.compute_language_confidence_values(corpus)
        if not confs:
            return None
        top = confs[0]  # sorted desc by lingua
        mapping = {
            Language.ENGLISH: "en",
            Language.JAPANESE: "ja",
            Language.CHINESE: "zh",
            Language.KOREAN: "ko",
        }
        return mapping.get(top.language)
    except Exception:
        logger.debug("lingua detection raised", exc_info=True)
        return None


def detect_slide_language(slides: list[SlideText]) -> str:
    """Detect the dominant language of extracted slide text.

    Returns one of: 'en', 'ja', 'zh', 'ko'. Defaults to 'en' if no
    alphabetic content was extracted.

    Strategy:
      1. If lingua-language-detector is installed, prefer it (more accurate
         on short / mixed text). Opt-in via `pip install lingua-language-detector`.
      2. Otherwise fall back to a zero-dep script-aware classifier:
         - Hiragana / Katakana never occur in Chinese or Korean text, so
           ANY meaningful kana presence → Japanese.
         - Hangul never occurs in Chinese or Japanese text, so dominant
           Hangul → Korean.
         - Han ideographs without kana → Chinese.

    Both paths handle real-world mixed decks (e.g. a Japanese business
    deck with English brand names like "Dell Technologies", numbered
    tables, and ASCII dates) which the prior CJK/total-chars ratio failed on.
    """
    via_lingua = _detect_via_lingua(slides)
    if via_lingua is not None:
        return via_lingua

    c = _script_counts(slides)
    cjk_total = c["hiragana"] + c["katakana"] + c["han"]
    alpha_total = c["latin"] + cjk_total + c["hangul"]
    if alpha_total == 0:
        return "en"

    # Hangul: Korean. Strong, script-exclusive signal.
    if c["hangul"] > 0 and c["hangul"] >= max(c["latin"], cjk_total) * 0.5:
        return "ko"

    # Kana presence is a hard "Japanese" signal — these scripts cannot
    # appear in Chinese or Korean text. Even a small amount in an
    # otherwise English-looking deck is decisive.
    kana = c["hiragana"] + c["katakana"]
    if kana >= 10 or (kana > 0 and kana / alpha_total >= 0.02):
        return "ja"

    # Han ideographs with no kana → Chinese.
    if c["han"] > 0 and c["han"] / alpha_total >= 0.10:
        return "zh"

    return "en"
