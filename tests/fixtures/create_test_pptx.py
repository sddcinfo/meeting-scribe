"""Generate a test PPTX fixture for slide translation tests.

Run once to create test_slides.pptx in this directory:
    python tests/fixtures/create_test_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).parent


def create_test_pptx() -> Path:
    """Create a small test PPTX with text boxes and a table."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide with two text boxes
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
    txBox1 = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    tf1 = txBox1.text_frame
    p1 = tf1.paragraphs[0]
    run1 = p1.add_run()
    run1.text = "Quarterly Revenue Report"
    run1.font.size = Pt(32)
    run1.font.bold = True

    txBox2 = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "APAC Region Overview"
    run2.font.size = Pt(24)

    # Slide 2: Content with table
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Regional Performance"
    run3.font.size = Pt(28)

    # Add a 3x2 table
    table_shape = slide2.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(6), Inches(3))
    table = table_shape.table

    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Growth"
    table.cell(1, 0).text = "Japan"
    table.cell(1, 1).text = "15%"
    table.cell(2, 0).text = "Australia"
    table.cell(2, 1).text = "23%"

    # Slide 3: Simple content
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox4 = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    tf4 = txBox4.text_frame
    p4 = tf4.paragraphs[0]
    run4 = p4.add_run()
    run4.text = "Thank you for your attention"
    run4.font.size = Pt(36)

    output = HERE / "test_slides.pptx"
    prs.save(str(output))
    return output


if __name__ == "__main__":
    path = create_test_pptx()
    print(f"Created: {path}")
