"""Tests for the slides conversion pipeline (convert.py).

These tests exercise validation, text extraction, reinsertion,
and the structured translation response validation.
"""

from __future__ import annotations

import asyncio
import tempfile
import zipfile
from pathlib import Path

import pytest

# Skip entire module if python-pptx is not installed
pptx = pytest.importorskip("pptx")

from meeting_scribe.slides.convert import (
    build_minimal_pptx,
    convert_pptx_to_images,
    detect_slide_language,
    extract_embedded_thumbnail,
    extract_text_from_pptx,
    reinsert_translated_text,
    render_first_slide_fast,
    render_partial_translated,
    slides_to_json,
    validate_pptx_contents,
)
from meeting_scribe.slides.job import (
    EXPRESS_BATCH_SCHEDULE,
    EXPRESS_BATCH_STEADY,
    SlideJobRunner,
    _DeckStats,
    express_batch_threshold,
)
from meeting_scribe.slides.models import SlideText, TextRun


def _slides_from_text(*texts: str) -> list[SlideText]:
    """Build SlideText fixtures from raw strings (one slide each)."""
    out: list[SlideText] = []
    for i, text in enumerate(texts):
        out.append(
            SlideText(
                index=i,
                runs=[
                    TextRun(
                        id=f"s{i}_r0",
                        slide_index=i,
                        shape_id=0,
                        para_index=0,
                        run_index=0,
                        text=text,
                    )
                ],
            )
        )
    return out


# ── Fixtures ─────────────────────────────────────────────────


@pytest.fixture
def test_pptx(tmp_path) -> Path:
    """Create a minimal test PPTX with text boxes and a table."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    run = txBox.text_frame.paragraphs[0].add_run()
    run.text = "Quarterly Report"
    run.font.size = Pt(32)
    run.font.bold = True

    txBox2 = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    run2 = txBox2.text_frame.paragraphs[0].add_run()
    run2.text = "APAC Overview"

    # Slide 2: Table
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    txBox3 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    run3 = txBox3.text_frame.paragraphs[0].add_run()
    run3.text = "Performance Data"

    table_shape = slide2.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Growth"
    table.cell(1, 0).text = "Japan"
    table.cell(1, 1).text = "15%"

    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def macro_pptx(tmp_path) -> Path:
    """Create a PPTX-like ZIP with a vbaProject.bin to test rejection."""
    path = tmp_path / "macro.pptx"
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("[Content_Types].xml", '<?xml version="1.0"?><Types></Types>')
        zf.writestr("ppt/vbaProject.bin", b"fake vba content")
    return path


@pytest.fixture
def oversized_entry_zip(tmp_path) -> Path:
    """Create a ZIP where one entry claims > 100 MB uncompressed (ZIP bomb check)."""
    path = tmp_path / "bomb.pptx"
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("[Content_Types].xml", '<?xml version="1.0"?><Types></Types>')
        # Write a small entry but we'll test the count/size checks separately
        zf.writestr("ppt/slides/slide1.xml", "x" * 100)
    return path


@pytest.fixture
def not_a_zip(tmp_path) -> Path:
    """Create a file that is not a ZIP."""
    path = tmp_path / "notazip.pptx"
    path.write_text("This is not a ZIP file")
    return path


# ── Validation tests ─────────────────────────────────────────


class TestValidation:
    def test_valid_pptx(self, test_pptx):
        result = validate_pptx_contents(test_pptx)
        assert result.valid
        assert result.slide_count == 2
        assert result.error is None

    def test_reject_not_zip(self, not_a_zip):
        result = validate_pptx_contents(not_a_zip)
        assert not result.valid
        assert "ZIP" in (result.error or "")

    def test_reject_not_pptx_content(self, macro_pptx):
        """ZIP without valid PPTX content is rejected."""
        result = validate_pptx_contents(macro_pptx)
        assert not result.valid


# ── Text extraction tests ────────────────────────────────────


class TestTextExtraction:
    def test_extract_text_boxes(self, test_pptx):
        slides = extract_text_from_pptx(test_pptx)
        assert len(slides) == 2

        # Slide 1 has two text boxes
        s1_texts = [r.text for r in slides[0].runs]
        assert "Quarterly Report" in s1_texts
        assert "APAC Overview" in s1_texts

    def test_extract_table_cells(self, test_pptx):
        slides = extract_text_from_pptx(test_pptx)

        # Slide 2 has a text box + table
        s2_texts = [r.text for r in slides[1].runs]
        assert "Performance Data" in s2_texts
        assert "Region" in s2_texts
        assert "Japan" in s2_texts
        assert "15%" in s2_texts

    def test_run_ids_are_unique(self, test_pptx):
        slides = extract_text_from_pptx(test_pptx)
        all_ids = [r.id for s in slides for r in s.runs]
        assert len(all_ids) == len(set(all_ids)), "Run IDs must be unique"

    def test_run_ids_contain_slide_index(self, test_pptx):
        slides = extract_text_from_pptx(test_pptx)
        for slide in slides:
            for run in slide.runs:
                assert run.id.startswith(f"s{slide.index}_")

    def test_serialization_roundtrip(self, test_pptx):
        slides = extract_text_from_pptx(test_pptx)
        data = slides_to_json(slides)
        assert isinstance(data, list)
        assert len(data) == 2
        for slide_data in data:
            assert "index" in slide_data
            assert "runs" in slide_data
            for run in slide_data["runs"]:
                assert "id" in run
                assert "text" in run


# ── Text reinsertion tests ───────────────────────────────────


class TestReinsertion:
    def test_reinsert_preserves_formatting(self, test_pptx, tmp_path):
        # Extract text
        slides = extract_text_from_pptx(test_pptx)
        first_run = slides[0].runs[0]

        # Create translations
        translations = [{"id": first_run.id, "translated": "Translated Title"}]

        # Reinsert
        output = tmp_path / "translated.pptx"
        reinsert_translated_text(test_pptx, translations, output)

        # Verify the text was replaced
        from pptx import Presentation

        prs = Presentation(str(output))
        slide1 = prs.slides[0]
        texts = []
        for shape in slide1.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            texts.append(run.text)

        assert "Translated Title" in texts
        assert "Quarterly Report" not in texts

    def test_reinsert_only_affects_targeted_runs(self, test_pptx, tmp_path):
        slides = extract_text_from_pptx(test_pptx)
        first_run = slides[0].runs[0]

        translations = [{"id": first_run.id, "translated": "Changed"}]
        output = tmp_path / "partial.pptx"
        reinsert_translated_text(test_pptx, translations, output)

        # The second text box should remain unchanged
        new_slides = extract_text_from_pptx(output)
        s1_texts = [r.text for r in new_slides[0].runs]
        assert "APAC Overview" in s1_texts


# ── Language detection ──────────────────────────────────────


class TestLanguageDetection:
    """Regression tests for the script-aware language classifier.

    The detector must handle real-world mixed decks where digits, punctuation,
    and English brand names would otherwise dilute CJK ratios. The previous
    implementation misclassified a 25-slide Japanese deck as "en" because
    `total_chars` included dates and ASCII numbers.
    """

    def test_pure_english(self):
        slides = _slides_from_text("Quarterly Report", "Performance Data 2025")
        assert detect_slide_language(slides) == "en"

    def test_pure_japanese(self):
        slides = _slides_from_text("売上報告書", "業績データ")
        assert detect_slide_language(slides) == "ja"

    def test_japanese_with_brand_names_and_dates(self):
        # The real failure case: a Japanese deck dominated by ASCII content
        # (brand names, role abbreviations, dates) but with names + titles
        # in Japanese. Previous detector misclassified this as "en".
        slides = _slides_from_text(
            "FY26 1H MUFG APS",
            "2025/3/12",
            "Dell Technologies",
            "平澤 賢太 (クライアント・エグゼクティブ)",
            "中尾 岳 (AE)",
            "山高 行宏 (PSA)",
            "加藤 秀明 (PSA)",
            "梶原 康男 (DCS)",
        )
        assert detect_slide_language(slides) == "ja"

    def test_japanese_with_minimal_kana_among_ascii_table(self):
        # Common pattern: a Japanese slide with a numeric table + a few
        # particles/headers. Kana is small in absolute terms but its presence
        # is a script-exclusive JA signal that beats EN noise.
        slides = _slides_from_text(
            "売上の推移",  # kanji + hiragana particle "の"
            "Q1 100% Q2 95% Q3 110% Q4 88%",  # large EN/numeric content
            "前年比",
        )
        assert detect_slide_language(slides) == "ja"

    def test_chinese_no_kana(self):
        slides = _slides_from_text("季度报告", "业绩数据 2025")
        assert detect_slide_language(slides) == "zh"

    def test_korean(self):
        slides = _slides_from_text("분기 보고서", "실적 데이터")
        assert detect_slide_language(slides) == "ko"

    def test_empty_or_numeric_only_defaults_en(self):
        # Slide 1 of the MUFG deck only had "2025/3/12" extracted — no
        # alphabetic content. Should not crash; defaults to en in the
        # script-aware fallback. When lingua is installed it may return
        # an arbitrary code on numeric-only input, so this test pins
        # the script-aware path.
        try:
            import lingua  # noqa: F401

            pytest.skip("lingua installed — script-aware fallback bypassed")
        except ImportError:
            pass
        slides = _slides_from_text("2025/3/12", "100%", "$1,234")
        assert detect_slide_language(slides) == "en"

    def test_english_with_one_kana_word(self):
        # A single hiragana word in an otherwise-English deck trips JA in
        # the script-aware fallback (kana presence is script-exclusive).
        # Lingua, when installed, weighs the dominant language and would
        # call this English — that's arguably more accurate for picking
        # translation direction (the deck mostly needs EN→JA, not the
        # other way). This test pins the script-aware fallback only.
        try:
            import lingua  # noqa: F401

            pytest.skip("lingua installed — uses dominant-language verdict")
        except ImportError:
            pass
        slides = _slides_from_text(
            "Introduction",
            "Performance metrics",
            "Q1 results",
            "あらまし",  # 4 hiragana — script-exclusive JA signal
            "Conclusion",
            "Thank you",
        )
        # With kana >= 4 and >= 2% of alpha → ja
        assert detect_slide_language(slides) == "ja"

    def test_handful_of_loanword_katakana_in_english_deck_stays_en(self):
        # An English deck with ONE katakana word (e.g. a borrowed proper
        # noun) and lots of English content — should NOT be misclassified
        # as JA. Threshold: kana < 10 AND kana < 2% of alpha.
        slides = [
            SlideText(
                index=0,
                runs=[
                    TextRun(
                        id="s0_r0",
                        slide_index=0,
                        shape_id=0,
                        para_index=0,
                        run_index=0,
                        text=(
                            "Introduction to our approach for solving this " * 20
                            + "ハイブリッド"  # 1 katakana word, 6 chars
                        ),
                    )
                ],
            )
        ]
        assert detect_slide_language(slides) == "en"


# ── Progressive rendering (requires LibreOffice + poppler) ──


def _libreoffice_available() -> bool:
    import shutil

    return bool(shutil.which("libreoffice") and shutil.which("pdftoppm") and shutil.which("pdfinfo"))


@pytest.mark.skipif(
    not _libreoffice_available(),
    reason="LibreOffice/poppler required for end-to-end rendering",
)
class TestProgressiveRender:
    """End-to-end tests for the priority-first PNG render pipeline."""

    @pytest.fixture
    def four_slide_pptx(self, tmp_path) -> Path:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        for i in range(4):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            run = box.text_frame.paragraphs[0].add_run()
            run.text = f"Slide {i + 1}"
            run.font.size = Pt(40)
            run.font.bold = True
        path = tmp_path / "four.pptx"
        prs.save(str(path))
        return path

    def test_renders_all_slides_with_progress_callback(
        self, four_slide_pptx, tmp_path
    ):
        out = tmp_path / "original"
        events: list[tuple[int, int]] = []

        count = convert_pptx_to_images(
            four_slide_pptx, out, progress_cb=lambda i, t: events.append((i, t))
        )

        assert count == 4
        assert len(events) == 4
        # Slide 1 must fire FIRST (priority pass)
        assert events[0][0] == 0, "slide 1 should be the first progress event"
        # Total in every event matches actual slide count
        assert all(e[1] == 4 for e in events)
        # Files on disk
        for i in range(1, 5):
            assert (out / f"slide_{i:03d}.png").exists()
        assert (out / "original.pdf").exists()

    def test_render_partial_translated_only_emits_targeted_slides(
        self, four_slide_pptx, tmp_path
    ):
        # Translate text on slides 0 and 1 only; render only those slides
        slides = extract_text_from_pptx(four_slide_pptx)
        translations = []
        for slide_idx in (0, 1):
            for run in slides[slide_idx].runs:
                translations.append({"id": run.id, "translated": f"TR-{slide_idx}"})

        out = tmp_path / "translated"
        work = tmp_path / "work"
        rendered = render_partial_translated(
            four_slide_pptx, translations, [0, 1], out, work
        )

        assert sorted(rendered) == [0, 1]
        assert (out / "slide_001.png").exists()
        assert (out / "slide_002.png").exists()
        # Slides 2 and 3 must NOT have been rendered by the express lane
        assert not (out / "slide_003.png").exists()
        assert not (out / "slide_004.png").exists()

    def test_render_first_slide_fast_produces_only_one_png(
        self, four_slide_pptx, tmp_path
    ):
        out = tmp_path / "slide_001.png"
        work = tmp_path / "work"
        ok = render_first_slide_fast(four_slide_pptx, out, work)
        assert ok is True
        assert out.exists()
        # No other slide PNGs leak into the work or output dir
        leaked = [p for p in tmp_path.glob("slide_*.png") if p.name != "slide_001.png"]
        assert leaked == []


class TestMinimalPptx:
    """Stripping a PPTX to one slide for fast LibreOffice render."""

    @pytest.fixture
    def six_slide_pptx(self, tmp_path) -> Path:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        for i in range(6):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            run = box.text_frame.paragraphs[0].add_run()
            run.text = f"Slide {i + 1}"
            run.font.size = Pt(40)
        path = tmp_path / "six.pptx"
        prs.save(str(path))
        return path

    def test_keeps_only_requested_slide(self, six_slide_pptx, tmp_path):
        from pptx import Presentation

        out = tmp_path / "minimal.pptx"
        kept = build_minimal_pptx(six_slide_pptx, [0], out)
        assert kept == 1
        prs = Presentation(str(out))
        assert len(prs.slides) == 1
        # The one kept slide is the first one
        texts = []
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            texts.append(run.text)
        assert "Slide 1" in texts

    def test_keeps_multiple_slides_in_order(self, six_slide_pptx, tmp_path):
        from pptx import Presentation

        out = tmp_path / "two.pptx"
        kept = build_minimal_pptx(six_slide_pptx, [0, 2], out)
        assert kept == 2
        prs = Presentation(str(out))
        assert len(prs.slides) == 2
        # Verify the two kept slides are slides 1 and 3 (not 1 and 2)
        all_texts = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text.strip():
                                all_texts.add(run.text)
        assert "Slide 1" in all_texts
        assert "Slide 3" in all_texts
        assert "Slide 2" not in all_texts


class TestEmbeddedThumbnail:
    """Reading docProps/thumbnail.* out of a PPTX zip."""

    def test_returns_false_when_no_thumbnail(self, tmp_path):
        # Bare zip with no docProps/thumbnail.* entries — common for decks
        # exported by some Google Slides flows.
        path = tmp_path / "nothumb.pptx"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("[Content_Types].xml", "<x/>")

        dest = tmp_path / "out.png"
        assert extract_embedded_thumbnail(path, dest) is False
        assert not dest.exists()

    def test_python_pptx_decks_carry_a_thumbnail(self, tmp_path):
        """Pleasant surprise: python-pptx writes a default thumbnail too.
        That means the instant placeholder path covers all four sources we
        commonly see (PowerPoint, Keynote, python-pptx, sometimes Google).
        """
        from pptx import Presentation

        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[5])
        path = tmp_path / "withthumb.pptx"
        prs.save(str(path))

        dest = tmp_path / "thumb.png"
        assert extract_embedded_thumbnail(path, dest) is True
        assert dest.exists()

    def test_extracts_when_present(self, tmp_path):
        # Build a fake PPTX-like zip with a known JPEG thumbnail entry
        # so we can verify extraction without depending on PowerPoint.
        from io import BytesIO

        from PIL import Image

        # Tiny 16x16 JPEG
        buf = BytesIO()
        Image.new("RGB", (16, 16), color=(120, 60, 60)).save(buf, format="JPEG")
        jpeg_bytes = buf.getvalue()

        path = tmp_path / "withthumb.pptx"
        with zipfile.ZipFile(path, "w") as zf:
            zf.writestr("[Content_Types].xml", "<x/>")
            zf.writestr("docProps/thumbnail.jpeg", jpeg_bytes)

        dest = tmp_path / "thumb.png"
        assert extract_embedded_thumbnail(path, dest) is True
        assert dest.exists()
        # Pillow re-encoded as PNG → file should start with the PNG magic
        assert dest.read_bytes()[:4] == b"\x89PNG"


# ── Parse-failure instrumentation (Phase 1 prereq #3) ────────────
#
# The plan's Phase 4a A/B can only detect the failure mode it's
# measuring if parse-failures are counted, not silently swallowed.
# These tests pin that contract: give `_translate_slide_text` a
# deliberately malformed LLM response, assert the counters move.
#
# Scoped narrowly to the parser — avoids spinning up a full slide
# pipeline.  See tests/test_drain.py for the pure SlideJobRunner
# harness pattern re-used here.


class TestSlideParseInstrumentation:
    """Malformed LLM responses must bump parse_failures, not vanish."""

    @staticmethod
    def _make_runner(mock_response: str | None) -> SlideJobRunner:
        """Build a SlideJobRunner whose translate_fn returns *mock_response*."""

        async def _mock_translate(text, source_lang, target_lang, system_prompt, max_tokens):
            return mock_response

        async def _noop_broadcast(_payload):
            pass

        # meetings_dir just has to exist; the parse tests never write there.
        tmp = tempfile.mkdtemp(prefix="scribe-parse-test-")
        return SlideJobRunner(
            meetings_dir=Path(tmp),
            translate_fn=_mock_translate,
            broadcast_fn=_noop_broadcast,
        )

    @staticmethod
    def _runs(*texts: str) -> list[dict]:
        """Build the runs list the parser expects (id + text per entry)."""
        return [{"id": f"r{i}", "text": t} for i, t in enumerate(texts)]

    @staticmethod
    def _register_stats(runner: SlideJobRunner, deck_id: str = "test-deck") -> _DeckStats:
        stats = _DeckStats()
        runner._deck_stats[deck_id] = stats
        runner.active_deck_id = deck_id
        return stats

    def test_regex_miss_bumps_parse_failures(self):
        # Model returns prose (one sentence per line) instead of the
        # numbered-list format we asked for.  Each line that doesn't
        # match `\d+[.)]\s*...` increments parse_failures individually.
        runner = self._make_runner("The first sentence.\nThe second sentence.")
        stats = self._register_stats(runner)

        results = asyncio.run(
            runner._translate_slide_text(
                self._runs("hello", "world"),
                source_lang="en",
                target_lang="ja",
            )
        )

        assert results == []
        assert stats.runs_requested == 2
        assert stats.runs_returned == 0
        # Two lines, both fail the `\d+[.)]` regex.
        assert stats.parse_failures == 2

    def test_empty_response_counts_every_run_as_parse_failure(self):
        runner = self._make_runner("")
        stats = self._register_stats(runner)

        results = asyncio.run(
            runner._translate_slide_text(
                self._runs("hello", "world", "foo"),
                source_lang="en",
                target_lang="ja",
            )
        )

        assert results == []
        assert stats.parse_failures == 3

    def test_none_response_counts_every_run_as_parse_failure(self):
        # Same code path as empty-string: translate_fn returned None
        # (backend error) → every run treated as a parse failure so
        # the stats log can distinguish "happy path" from "silent skip".
        runner = self._make_runner(None)
        stats = self._register_stats(runner)

        results = asyncio.run(
            runner._translate_slide_text(
                self._runs("hello"),
                source_lang="en",
                target_lang="ja",
            )
        )

        assert results == []
        assert stats.parse_failures == 1

    def test_out_of_range_id_bumps_parse_failures(self):
        # Model hallucinates "3." when only 2 runs were requested.
        runner = self._make_runner("1. one\n2. two\n3. three")
        stats = self._register_stats(runner)

        results = asyncio.run(
            runner._translate_slide_text(
                self._runs("hello", "world"),
                source_lang="en",
                target_lang="ja",
            )
        )

        assert len(results) == 2
        assert stats.runs_requested == 2
        assert stats.runs_returned == 2
        # Index 3 (1-indexed) was out of range → counted.
        assert stats.parse_failures == 1

    def test_duplicate_id_bumps_id_coverage_failures(self):
        # Model emits the same ID twice; the second occurrence counts
        # against ID coverage so Phase 4b's JSON-schema path has a
        # direct apples-to-apples failure metric.
        runner = self._make_runner("1. one\n1. also one\n2. two")
        stats = self._register_stats(runner)

        results = asyncio.run(
            runner._translate_slide_text(
                self._runs("hello", "world"),
                source_lang="en",
                target_lang="ja",
            )
        )

        assert len(results) == 2
        assert stats.id_coverage_failures == 1
        assert stats.parse_failures == 0

    def test_missing_id_bumps_id_coverage_failures(self):
        # Model returns only 1 of 3 requested IDs → 2 missing.
        runner = self._make_runner("1. one")
        stats = self._register_stats(runner)

        results = asyncio.run(
            runner._translate_slide_text(
                self._runs("a", "b", "c"),
                source_lang="en",
                target_lang="ja",
            )
        )

        assert len(results) == 1
        assert stats.id_coverage_failures == 2

    def test_clean_response_leaves_counters_at_zero_failures(self):
        runner = self._make_runner("1. ichi\n2. ni\n3. san")
        stats = self._register_stats(runner)

        results = asyncio.run(
            runner._translate_slide_text(
                self._runs("one", "two", "three"),
                source_lang="en",
                target_lang="ja",
            )
        )

        assert len(results) == 3
        assert stats.parse_failures == 0
        assert stats.id_coverage_failures == 0
        assert stats.runs_returned == 3


# ── Regression: concurrent pdftoppm tmp-file collision ──────────


class TestPdftoppmTmpFileIsolation:
    """Regression guard for the slide-switch bug.

    Two render paths target the same final dest (``original/slide_001.png``):
    the express-first-slide renderer and the bulk-originals renderer.
    Before the fix, both used a shared ``slide_001.tmp.png`` tmp file,
    so a race between them could leave the bulk renderer looking at
    an orphan inode, raising ``pdftoppm produced no output`` and
    aborting the bulk render after slide 1.  That produced the
    on-disk state ``original/slide_001.png + original.pdf only``,
    which the user experienced as "first slide works, next slide
    breaks".

    This test drives ``_pdftoppm_single_page`` concurrently for the
    same dest and asserts both callers complete cleanly.
    """

    def test_concurrent_same_dest_does_not_orphan_tmp_files(self, tmp_path):
        """Simulates the express + bulk race using a stub pdftoppm."""
        from meeting_scribe.slides import convert as convert_mod

        # Capture every temp prefix pdftoppm is asked to write to.
        tmp_prefixes_seen: list[str] = []

        def _fake_subprocess_run(argv, *args, **kwargs):
            # pdftoppm cmdline layout: ``[..., "-singlefile", pdf, prefix]``.
            # ``prefix`` is the last positional.  Write a PNG there as if
            # pdftoppm had produced it.
            prefix = Path(argv[-1])
            out = Path(str(prefix) + ".png")
            out.write_bytes(b"\x89PNG\r\n\x1a\nstub")
            tmp_prefixes_seen.append(out.name)

            class _Result:
                returncode = 0
                stdout = ""
                stderr = ""

            return _Result()

        import subprocess as _sub

        orig_run = _sub.run
        _sub.run = _fake_subprocess_run
        try:
            # Use a stub PDF file — the stub run never reads it, so
            # an empty path is fine.
            pdf = tmp_path / "stub.pdf"
            pdf.write_bytes(b"%PDF-stub")

            dest = tmp_path / "slide_001.png"

            # Run two calls sequentially but NOTE the test's value is
            # in the *tmp prefix* being unique per invocation — which
            # we can verify directly from the captured list.  The real
            # race requires threads; isolation of the tmp name is the
            # invariant.
            convert_mod._pdftoppm_single_page(pdf, 1, dest)
            assert dest.exists()
            dest.unlink()

            convert_mod._pdftoppm_single_page(pdf, 1, dest)
            assert dest.exists()

            # Two invocations must produce two distinct temp prefix
            # filenames, otherwise concurrent callers could race on the
            # same file.
            assert len(tmp_prefixes_seen) == 2
            assert tmp_prefixes_seen[0] != tmp_prefixes_seen[1], (
                "pdftoppm tmp prefix must be unique per invocation to "
                "prevent the express+bulk render collision"
            )
            # Both tmp filenames should START with the dest stem
            # but carry a distinct randomizer.
            for name in tmp_prefixes_seen:
                assert name.startswith("slide_001."), name
                assert name.endswith(".tmp.png"), name
        finally:
            _sub.run = orig_run


class TestExpressBatchSchedule:
    """Size-adaptive batching gives the user an early first-paint: the
    first three batches flush at 1, 2, 3 slides respectively, then
    steady-state flushes at 6. Protects the UX fix for the "all or
    nothing at 6 slides" complaint.
    """

    def test_schedule_starts_small_then_steady(self) -> None:
        assert EXPRESS_BATCH_SCHEDULE == (1, 2, 3)
        assert EXPRESS_BATCH_STEADY == 6

    def test_threshold_per_batch_index(self) -> None:
        assert express_batch_threshold(0) == 1
        assert express_batch_threshold(1) == 2
        assert express_batch_threshold(2) == 3
        assert express_batch_threshold(3) == EXPRESS_BATCH_STEADY
        assert express_batch_threshold(4) == EXPRESS_BATCH_STEADY
        assert express_batch_threshold(42) == EXPRESS_BATCH_STEADY

    def test_first_paint_threshold_is_one(self) -> None:
        """The first batch MUST fire at one buffered slide so the user
        sees a translation render within one LO-invocation of the first
        slide finishing translation, not within 6."""
        assert express_batch_threshold(0) == 1

    def test_fifty_slide_deck_batch_count_matches_schedule(self) -> None:
        """For a 50-slide deck, verify the batch schedule: 1+2+3 = 6
        slides across the first three batches, then (50-6)/6 = 7.33
        → 8 more batches (final is partial). Total: 11 batches."""
        remaining = 50
        batches = []
        batches_fired = 0
        while remaining > 0:
            threshold = express_batch_threshold(batches_fired)
            take = min(threshold, remaining)
            batches.append(take)
            remaining -= take
            batches_fired += 1
        assert batches[:3] == [1, 2, 3]
        assert all(b == 6 or b < 6 for b in batches[3:])
        assert sum(batches) == 50
        # Old flat-6 schedule: ceil(50/6) = 9 batches. New: 11.
        # Acceptable trade for progressive UX.
        assert len(batches) == 11


class TestPartialTranslatedWorkDirIsolation:
    """Concurrent express-batch renders must use disjoint work dirs.

    A shared `_express` dir would race: every call writes
    source.pptx, partial_translated.pptx, and express_minimal.pptx to
    the same paths. Under SCRIBE_SLIDE_RENDER_PARALLELISM=4 that would
    trample intermediate files and produce corrupted PDFs. The uuid
    suffix on the per-invocation scratch dir is the guard.
    """

    def test_work_dir_has_uuid_suffix(self, monkeypatch, tmp_path):
        import asyncio

        from meeting_scribe.slides import worker as worker_mod

        seen_work_dirs: list[Path] = []

        def _spy_render(
            source_pptx: Path,
            translations: list[dict],
            slide_indices_0based: list[int],
            output_dir: Path,
            work_dir: Path,
        ) -> list[int]:
            seen_work_dirs.append(work_dir)
            return []

        monkeypatch.setattr(
            "meeting_scribe.slides.convert.render_partial_translated",
            _spy_render,
        )

        deck_dir = tmp_path / "deck"
        deck_dir.mkdir()
        translated = deck_dir / "translated"
        translated.mkdir()

        async def _run_two() -> None:
            await asyncio.gather(
                worker_mod.run_partial_translated_render(
                    b"fake-pptx-bytes-1", [], [0], translated
                ),
                worker_mod.run_partial_translated_render(
                    b"fake-pptx-bytes-2", [], [1], translated
                ),
            )

        asyncio.run(_run_two())

        assert len(seen_work_dirs) == 2
        # Both work dirs must be DISTINCT (otherwise race).
        assert seen_work_dirs[0] != seen_work_dirs[1], (
            "parallel run_partial_translated_render calls must use "
            "disjoint work dirs to avoid source.pptx / minimal.pptx "
            "file collisions"
        )
        # Both must carry the `_express_` prefix + uuid-hex suffix.
        for wd in seen_work_dirs:
            assert wd.name.startswith("_express_"), wd
            suffix = wd.name[len("_express_"):]
            assert len(suffix) == 32, f"expected uuid4 hex (32 chars), got {suffix!r}"


class TestAtomicWriteJsonRace:
    """meta.tmp was a fixed filename; concurrent writers raced on rename.

    The bug manifested as a silent translation skip when two threads
    raced on meta.json during the extract-text + render-originals
    parallel phase (2026-04-20). Guard: uuid-suffixed tmp filenames.
    """

    def test_concurrent_writes_to_same_path_succeed(self, tmp_path):
        import threading

        from meeting_scribe.slides.worker import _atomic_write_json

        target = tmp_path / "meta.json"
        errors: list[Exception] = []

        def _writer(i: int) -> None:
            try:
                for _ in range(5):
                    _atomic_write_json(target, {"worker": i})
            except Exception as exc:
                errors.append(exc)

        threads = [threading.Thread(target=_writer, args=(i,)) for i in range(8)]
        for t in threads:
            t.start()
        for t in threads:
            t.join()

        assert not errors, f"concurrent _atomic_write_json raised: {errors}"
        # Target file must have the final shape — content from one of the
        # writers — not be missing or partial.
        import json as _json

        data = _json.loads(target.read_text())
        assert "worker" in data

